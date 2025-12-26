from pptx import Presentation
from summarize import get_summarized_news
from datetime import datetime


def create_ppt(news_list):
    prs = Presentation("AIWeeklyReport_format.pptx")
    slide = prs.slides[0]

    # Number of report and date
    issue_num = input("Number of the report (ex: 14): ")
    input_date = input("Date (ex: 2025년 12월 19일, today date(enter): ")
    if not input_date:
        date_str = datetime.now().strftime("%Y년 %m월 %d일")
    else:
        date_str = input_date
    slide.shapes[4].text = f"제 {issue_num}호 | {date_str}"

    for i, shape in enumerate(slide.shapes):
        text = shape.text if shape.has_text_frame else "no text"
        print(f"INDEX: {i} | NAME: {shape.name} | CONTENT: '{text}'")

    # for news in news_list:
    #     slide = prs.slides.add_slide(prs.slide_layouts[1])
    #     slide.shapes.title.text = news["title"]
    #     slide.placeholders[1].text = news["summary"]

    prs.save("insurance_ai_news_summary.pptx")

if __name__ == "__main__":
    news = get_summarized_news()
    print(news)
    create_ppt(news)