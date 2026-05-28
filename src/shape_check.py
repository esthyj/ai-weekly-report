"""shape_check.py — templates/AIWeeklyReport_format.pptx 각 shape에 작성된 내용 출력.

실행:
    python -m src.shape_check
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from config import PPT_TEMPLATE_FILE


def walk_shapes(shapes, slide_idx, prefix=""):
    for shape_idx, shape in enumerate(shapes):
        print(f"[슬라이드 {slide_idx} / {prefix}shape {shape_idx}] "
              f"{shape.name} (type={shape.shape_type})")

        # 그룹이면 안으로 재귀
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk_shapes(shape.shapes, slide_idx, prefix=f"{prefix}{shape_idx}-")
            continue

        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text_frame.text
            if text.strip():
                for line in text.splitlines():
                    print(f"  {line}")
        print()


def main() -> None:
    prs = Presentation(str(PPT_TEMPLATE_FILE))
    for slide_idx, slide in enumerate(prs.slides):
        walk_shapes(slide.shapes, slide_idx)


if __name__ == "__main__":
    main()