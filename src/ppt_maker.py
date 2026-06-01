import re
import math
from pathlib import Path
from typing import Union, Sequence
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from .config import PPT_TEMPLATE_FILE

ShapePath = Union[int, Sequence[int]]

# ============================================================
# Settings
# ============================================================
TAG_RE = re.compile(r'\[(Title|Summary\d*|Insight)\]\s*', re.IGNORECASE)

# 인라인 서식 마크업 — 웹 final-review 에디터에서 직렬화한 형태
#   {b}…{/b}  {i}…{/i}  {u}…{/u}  {size=14}…{/size}  {color=#c00000}…{/color}
# 열림 태그는 value 포함, 닫힘 태그는 value 없음을 모두 매칭.
INLINE_RE = re.compile(
    r"\{(/?)(b|i|u|size(?:=\d+)?|color(?:=#[0-9a-fA-F]{6})?)\}"
)

# 태그별 스타일 설정 (prefix, font_name, font_size, underline, split_lines)
TAG_STYLES = {
    "title":   ("",   "한화고딕 B",  12, False, False),
    "summary": ("• ", "한화고딕 EL", 12, False, True),
    "insight": ("➔ ", "한화고딕 B",  12, True,  True),
}
DEFAULT_STYLE = ("", "한화고딕 EL", 12, False, True)

# Slide 0 shape indices — matches templates/AIWeeklyReport_format.pptx structure.
# If the template is restructured, re-check with list_all_shapes() and update here.
META_SHAPE_INDEX: ShapePath = 4         # 호수 + 날짜
NEWS_SHAPE_INDEX: ShapePath = (10, 0)   # 그룹 10 안의 자식 0 — 뉴스 요약 본문
AILAB_SHAPE_INDEX: ShapePath = (9, 0)   # 그룹 9 안의 자식 0 — AI Lab 요약 본문


# ============================================================
# Utility Functions
# ============================================================

# 태그에 맞는 스타일 반환 (Summary1, Summary2 등 모두 summary 스타일 적용)
def get_tag_style(tag: str):
    tag = tag.lower()
    if tag.startswith("summary"):
        return TAG_STYLES["summary"]
    return TAG_STYLES.get(tag, DEFAULT_STYLE)


# Split a long text by tag (header) and extract each tag's section content into a list
def parse_sections(text: str):
    matches = list(TAG_RE.finditer(text))
    return [
        (m.group(1).lower(), text[m.end():matches[i+1].start() if i+1 < len(matches) else len(text)].strip())
        for i, m in enumerate(matches)
        if text[m.end():matches[i+1].start() if i+1 < len(matches) else len(text)].strip()
    ]


# Return specific shape. shape_index can be:
#   - int            : top-level shape index on the slide
#   - sequence[int]  : path descending into group shapes (e.g. (10, 0) → slide.shapes[10].shapes[0])
def find_shape_by_index(prs: Presentation, shape_index: ShapePath, slide_index: int = 0):
    if slide_index >= len(prs.slides):
        return None, None

    slide = prs.slides[slide_index]
    path = (shape_index,) if isinstance(shape_index, int) else tuple(shape_index)
    if not path:
        return None, None

    shapes = list(slide.shapes)
    shape = None
    for depth, idx in enumerate(path):
        if idx >= len(shapes):
            return None, None
        shape = shapes[idx]
        if depth < len(path) - 1:
            # Need to descend — current shape must be a group
            if not hasattr(shape, "shapes"):
                return None, None
            shapes = list(shape.shapes)

    return slide, shape


# Add a styled text run
def add_styled_run(paragraph, text, font_name, font_size, underline=False, color=None):
    r = paragraph.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.underline = underline
    if color:
        r.font.color.rgb = color


# Parse inline markup into a list of (text_segment, overrides) pairs.
# overrides may contain keys: bold, italic, underline, size (int), color ('#rrggbb').
# 마크업이 전혀 없으면 [(text, {})] 한 개 — 기존 평문 입력과 결과 동일.
def parse_inline(text: str):
    segments = []
    stack = []  # [("b", True) | ("i", True) | ("u", True) | ("size", 14) | ("color", "#c00000")]

    def current_overrides():
        ov = {}
        for typ, val in stack:
            if typ == "b":
                ov["bold"] = True
            elif typ == "i":
                ov["italic"] = True
            elif typ == "u":
                ov["underline"] = True
            elif typ == "size":
                ov["size"] = val
            elif typ == "color":
                ov["color"] = val
        return ov

    pos = 0
    for m in INLINE_RE.finditer(text):
        if m.start() > pos:
            segments.append((text[pos:m.start()], current_overrides()))

        closing = m.group(1) == "/"
        body = m.group(2)

        if body in ("b", "i", "u"):
            typ, val = body, True
        elif body == "size":  # 닫힘 태그: 가장 최근 size를 닫음
            typ, val = "size", None
        elif body.startswith("size="):
            typ, val = "size", int(body[len("size="):])
        elif body == "color":
            typ, val = "color", None
        elif body.startswith("color=#"):
            typ, val = "color", body[len("color="):]
        else:
            typ, val = None, None  # 정규식상 도달 불가

        if closing:
            # 같은 타입 중 가장 최근 항목을 제거 (관대 파싱)
            for i in range(len(stack) - 1, -1, -1):
                if stack[i][0] == typ:
                    stack.pop(i)
                    break
        else:
            stack.append((typ, val))

        pos = m.end()

    if pos < len(text):
        segments.append((text[pos:], current_overrides()))

    return segments


# Add one run with optional per-run style overrides on top of tag base style.
def add_run_with_overrides(paragraph, text, font_name, font_size, base_underline=False, overrides=None):
    if not text:
        return
    r = paragraph.add_run()
    r.text = text
    r.font.name = font_name
    ov = overrides or {}
    r.font.size = Pt(ov.get("size", font_size))
    r.font.underline = ov.get("underline", base_underline)
    r.font.italic = ov.get("italic", False)
    if ov.get("bold"):
        r.font.bold = True
    if "color" in ov:
        r.font.color.rgb = RGBColor.from_string(ov["color"].lstrip("#"))


# ============================================================
# Dynamic shape height
# ============================================================

EMU_PER_PT = 12700  # 1pt = 12700 EMU

# 한 글자의 가로 폭을 폰트 크기(em) 기준 단위로 환산.
# 한글/한자/전각 문자는 1.0em, 그 외(영문·숫자·기호·공백)는 약 0.55em 로 근사.
def _char_width_em(ch: str) -> float:
    o = ord(ch)
    if (0xAC00 <= o <= 0xD7A3       # 한글 음절
            or 0x3130 <= o <= 0x318F    # 한글 자모
            or 0x4E00 <= o <= 0x9FFF    # 한자
            or 0x3000 <= o <= 0x303F    # CJK 문장부호
            or 0xFF00 <= o <= 0xFFEF):  # 전각 영숫자/기호
        return 1.0
    return 0.55


# 채워진 텍스트 프레임이 실제로 차지할 세로 높이를 추정(EMU).
# word_wrap=True 기준으로 박스 폭에 맞춰 줄바꿈되는 줄 수를 문자 폭 합으로 근사한다.
# line_factor: 폰트 크기 대비 줄 높이 배율(단일 줄간격 ≈ 1.2~1.25).
def _estimate_text_frame_height(shape, line_factor: float = 1.25) -> int:
    tf = shape.text_frame
    avail_pt = (shape.width - tf.margin_left - tf.margin_right) / EMU_PER_PT
    if avail_pt <= 0:
        avail_pt = shape.width / EMU_PER_PT

    total_pt = 0.0
    for para in tf.paragraphs:
        sizes = [r.font.size.pt for r in para.runs if r.font.size is not None]
        font_pt = max(sizes) if sizes else 12.0
        score = sum(_char_width_em(ch) for r in para.runs for ch in r.text)

        units_per_line = max(avail_pt / font_pt, 1.0)
        lines = max(1, math.ceil(score / units_per_line)) if score else 1
        total_pt += lines * font_pt * line_factor

    return int(total_pt * EMU_PER_PT)


# 도형 높이를 텍스트 양에 맞춰 동적으로 조절한다.
# 본문 정렬이 MIDDLE 이므로 위·아래로 대칭적인 약간의 여유(padding)가 생긴다.
# (그룹 자식이라도 chExt==ext 라 cy 가 곧 렌더링 높이.)
def autosize_shape_height(shape, padding=Pt(3)):
    tf = shape.text_frame
    content = _estimate_text_frame_height(shape)
    shape.height = int(content + tf.margin_top + tf.margin_bottom + padding)


# ============================================================
# Write PPT
# ============================================================

# Add report number and date
def set_number_and_date(prs: Presentation, number: str, date: str,
                        shape_index: ShapePath = META_SHAPE_INDEX, slide_index: int = 0):
    """숫자와 날짜를 특정 TextBox에 입력"""
    _, shape = find_shape_by_index(prs, shape_index, slide_index)
    
    if not shape:
        raise ValueError(f'슬라이드 {slide_index}의 {shape_index}번째 shape을 찾지 못했습니다.')
    if not shape.has_text_frame:
        raise ValueError(f'{shape_index}번째 shape에 text_frame이 없습니다.')

    tf = shape.text_frame
    tf.clear()
    
    combined_text = f"제{number}호 | {date}"
    
    p = tf.paragraphs[0]
    add_styled_run(p, combined_text, "한화고딕 L", 11, color=RGBColor(0x6C, 0x6A, 0x67))


# Insert summarized text structured with tag-specific styles
def set_textbox_from_summarizedtxt(prs: Presentation, text: str,
                                    shape_index: ShapePath = NEWS_SHAPE_INDEX, slide_index: int = 0):
    # Find specific index shape
    _, shape = find_shape_by_index(prs, shape_index, slide_index)
    
    if not shape:
        raise ValueError(f'슬라이드 {slide_index}의 {shape_index}번째 shape을 찾지 못했습니다.')
    if not shape.has_text_frame:
        raise ValueError(f'{shape_index}번째 shape에 text_frame이 없습니다.')

    # Clear existing text frame
    tf = shape.text_frame
    tf.clear()

    # Return (tag, content) list
    sections = parse_sections(text)

    if not sections:
        add_styled_run(tf.paragraphs[0], text.strip(), "한화고딕 EL", 12)
        autosize_shape_height(shape)
        return

    first_para_used = False
    for tag, content in sections:
        # Find the style for the tag
        prefix, font_name, font_size, underline, split = get_tag_style(tag)
        lines = [ln.strip() for ln in content.splitlines() if ln.strip()] if split else [content.strip()]

        for line in filter(None, lines):
            p = tf.paragraphs[0] if not first_para_used and not tf.paragraphs[0].text else tf.add_paragraph()
            first_para_used = True

            # prefix(•, ➔)는 인라인 서식과 무관하게 항상 태그 기본 스타일로 출력
            if prefix:
                add_styled_run(p, prefix, font_name, font_size, underline)

            # 본문은 인라인 마크업을 파싱해 run 분할 — 마크업 없으면 한 개 run
            for seg, overrides in parse_inline(line):
                add_run_with_overrides(p, seg, font_name, font_size, underline, overrides)

        if tag == "insight":
            add_styled_run(tf.add_paragraph(), " ", "한화고딕 EL", 9)

    autosize_shape_height(shape)


# ============================================================
# Main Function
# ============================================================
# Create Report PPTX
def create_report(pptx_in: str, pptx_out: str, number: str, date: str,
                  text1: str, text2: str):

    # Check if template file exists
    if not Path(pptx_in).exists():
        raise FileNotFoundError(f"❌ PPT 템플릿 파일을 찾을 수 없습니다: {pptx_in}")

    prs = Presentation(pptx_in)
    
    # Step 1: Enter number of the report and date.
    set_number_and_date(prs, number, date, shape_index=META_SHAPE_INDEX, slide_index=0)

    # Step 2: Enter first summary text
    set_textbox_from_summarizedtxt(prs, text1, shape_index=NEWS_SHAPE_INDEX, slide_index=0)

    # Step 3: Enter second summary text
    set_textbox_from_summarizedtxt(prs, text2, shape_index=AILAB_SHAPE_INDEX, slide_index=0)
    
    # Save
    prs.save(pptx_out)
    print(f"  💾 {pptx_out} 저장 완료!")


# For debugging: output shape information for all slides
def list_all_shapes(pptx_path: str):
    prs = Presentation(pptx_path)
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n=== 슬라이드 {slide_idx} ===")
        for i, shape in enumerate(slide.shapes):
            name = getattr(shape, "name", "N/A")
            has_tf = hasattr(shape, "has_text_frame") and shape.has_text_frame
            text_preview = ""
            if has_tf and shape.text_frame.text:
                text_preview = shape.text_frame.text[:30].replace('\n', ' ') + "..."
            print(f"  [{i}] {name} (text_frame: {has_tf}) {text_preview}")


# Test (If needed)
if __name__ == "__main__":

    list_all_shapes(str(PPT_TEMPLATE_FILE))

    test_text1 = '''[Title] 테스트 제목 [Summary1] 요약1 내용 [Summary2] 요약2 내용 [Summary3] 요약3 내용 [Insight] 인사이트 내용'''
    test_text2 = '''[Title] AI Lab 테스트 [Summary1] AI Lab 요약1 [Summary2] AI Lab 요약2 [Insight] AI Lab 인사이트'''

    create_report(
        pptx_in=str(PPT_TEMPLATE_FILE),
        pptx_out="test_output.pptx",
        number="테스트",
        date="2025년 1월 1일",
        text1=test_text1,
        text2=test_text2
    )