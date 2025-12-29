import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ============================================================
# 설정
# ============================================================
TAG_RE = re.compile(r'\[(Title|Summary1|Summary2|Insight)\]\s*', re.IGNORECASE)

# 태그별 스타일 설정 (prefix, font_name, font_size, underline, split_lines)
TAG_STYLES = {
    "title":    ("",   "한화고딕 B",  12, False, False),
    "summary1": ("• ", "한화고딕 EL", 12, False, True),
    "summary2": ("• ", "한화고딕 EL", 12, False, True),
    "insight":  ("➔ ", "한화고딕 B",  12, True,  True),
}
DEFAULT_STYLE = ("", "한화고딕 EL", 12, False, True)


# ============================================================
# 유틸리티 함수
# ============================================================
def parse_sections(text: str):
    """텍스트에서 태그별 섹션 추출"""
    matches = list(TAG_RE.finditer(text))
    return [
        (m.group(1).lower(), text[m.end():matches[i+1].start() if i+1 < len(matches) else len(text)].strip())
        for i, m in enumerate(matches)
        if text[m.end():matches[i+1].start() if i+1 < len(matches) else len(text)].strip()
    ]


def find_shape_by_index(prs: Presentation, shape_index: int, slide_index: int = 0):
    """특정 슬라이드의 특정 인덱스 shape 반환"""
    if slide_index >= len(prs.slides):
        return None, None
    
    slide = prs.slides[slide_index]
    shapes = list(slide.shapes)
    
    if shape_index >= len(shapes):
        return None, None
    
    return slide, shapes[shape_index]


def add_styled_run(paragraph, text, font_name, font_size, underline=False, color=None):
    """스타일이 적용된 텍스트 run 추가"""
    r = paragraph.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.underline = underline
    if color:
        r.font.color.rgb = color


# ============================================================
# PPT 조작 함수
# ============================================================
def set_number_and_date(prs: Presentation, number: str, date: str, 
                        shape_index: int = 4, slide_index: int = 0):
    """숫자와 날짜를 특정 TextBox에 입력"""
    _, shape = find_shape_by_index(prs, shape_index, slide_index)
    
    if not shape:
        raise ValueError(f'슬라이드 {slide_index}의 {shape_index}번째 shape을 찾지 못했습니다.')
    if not shape.has_text_frame:
        raise ValueError(f'{shape_index}번째 shape에 text_frame이 없습니다.')

    tf = shape.text_frame
    tf.clear()
    
    combined_text = f"제{number}호 | {date}"
    
    p = tf.paragraphs[0]
    add_styled_run(p, combined_text, "한화고딕 L", 11, color=RGBColor(0x6C, 0x6A, 0x67))


def set_textbox_from_summarizedtxt(prs: Presentation, text: str, 
                                    shape_index: int = 15, slide_index: int = 0):
    """특정 슬라이드의 특정 인덱스 shape에 요약 텍스트 삽입"""
    _, shape = find_shape_by_index(prs, shape_index, slide_index)
    
    if not shape:
        raise ValueError(f'슬라이드 {slide_index}의 {shape_index}번째 shape을 찾지 못했습니다.')
    if not shape.has_text_frame:
        raise ValueError(f'{shape_index}번째 shape에 text_frame이 없습니다.')

    tf = shape.text_frame
    tf.clear()
    sections = parse_sections(text)

    if not sections:
        add_styled_run(tf.paragraphs[0], text.strip(), "한화고딕 EL", 12)
        return

    first_para_used = False
    for tag, content in sections:
        prefix, font_name, font_size, underline, split = TAG_STYLES.get(tag, DEFAULT_STYLE)
        lines = [ln.strip() for ln in content.splitlines() if ln.strip()] if split else [content.strip()]

        for line in filter(None, lines):
            p = tf.paragraphs[0] if not first_para_used and not tf.paragraphs[0].text else tf.add_paragraph()
            first_para_used = True
            add_styled_run(p, f"{prefix}{line}" if prefix else line, font_name, font_size, underline)

        if tag == "insight":
            add_styled_run(tf.add_paragraph(), " ", "한화고딕 EL", 9)


# ============================================================
# 메인 함수
# ============================================================
def create_report(pptx_in: str, pptx_out: str, number: str, date: str, 
                  text1: str, text2: str):
    """PPT 보고서 생성"""
    prs = Presentation(pptx_in)
    
    # 1단계: 숫자와 날짜 입력
    set_number_and_date(prs, number, date, shape_index=4, slide_index=0)
    
    # 2단계: 첫 번째 요약 텍스트 입력
    set_textbox_from_summarizedtxt(prs, text1, shape_index=15, slide_index=0)
    
    # 3단계: 두 번째 요약 텍스트 입력
    set_textbox_from_summarizedtxt(prs, text2, shape_index=16, slide_index=0)
    
    # 저장
    prs.save(pptx_out)
    print(f"  💾 {pptx_out} 저장 완료!")


def list_all_shapes(pptx_path: str):
    """디버깅용: 모든 슬라이드의 shape 정보 출력"""
    prs = Presentation(pptx_path)
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n=== 슬라이드 {slide_idx} ===")
        for i, shape in enumerate(slide.shapes):
            name = getattr(shape, "name", "N/A")
            has_tf = hasattr(shape, "has_text_frame") and shape.has_text_frame
            text_preview = ""
            if has_tf and shape.text_frame.text:
                text_preview = shape.text_frame.text[:30].replace('\n', ' ') + "..."
            print(f"  [{i}] {name} (text_frame: {has_tf}) {text_preview}")


# 테스트용 (직접 실행 시)
if __name__ == "__main__":
    # 디버깅: shape 목록 확인
    # list_all_shapes("AIWeeklyReport_format.pptx")
    
    # 테스트 텍스트로 보고서 생성
    test_text1 = '''[Title] 테스트 제목 [Summary1] 요약1 내용 [Summary2] 요약2 내용 [Insight] 인사이트 내용'''
    test_text2 = '''[Title] AI Lab 테스트 [Summary1] AI Lab 요약1 [Summary2] AI Lab 요약2 [Insight] AI Lab 인사이트'''
    
    create_report(
        pptx_in="AIWeeklyReport_format.pptx",
        pptx_out="test_output.pptx",
        number="테스트",
        date="2025년 1월 1일",
        text1=test_text1,
        text2=test_text2
    )
