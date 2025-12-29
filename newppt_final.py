import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor  
from summarize import get_summarized_news
from ailab_summarize import ailab_summarized  # ailab_summarize에서 import

summarized_text = get_summarized_news()
summarized_text2 = ailab_summarized()  # ailab_summarize.py의 return 값

# 테스트용
# summarized_text = '''
# [Title] 미래에셋생명, AI 기반 맞춤형 건광관리 서비스 '헬스케어 AI' 제공 [Summary1] 사용자의 건강검진 기록과 의료 데이터를 종합해 주요 질환의 발병 가능성을 예측하여, 사용자가 선제적으로 위험 인자를 관리할 수 있도록 도움 제공 [Summary2] 미래에셋생명의 AI 헬스케어 서비스는 질병 예측 AI, 기대 수명 예측, 의료비 예측, 개인 맞춤 건강 가이드 네 가지 핵심 특징을 갖추고 있음 [Insight] 당사도 향후 여성 헬스케어 서비스를 제공할 경우, 여성에게서 빈번하게 발생하는 주요 질환의 발병 가능성과 예상 의료비 기반으로 개인별 건강·재무 계획 수립을 지원할 수 있을것으로 기대됨
# '''

# summarized_text2 = '''
# [Title] AI Lab 관련 뉴스 제목 [Summary1] AI Lab 뉴스 요약1 [Summary2] AI Lab 뉴스 요약2 [Insight] AI Lab 뉴스 인사이트
# '''

TAG_RE = re.compile(r'\[(Title|Summary1|Summary2|Insight)\]\s*', re.IGNORECASE)

# 태그별 스타일 설정 (prefix, font_name, font_size, underline, split_lines)
TAG_STYLES = {
    "title":    ("",   "한화고딕 B",  12, False, False),
    "summary1": ("• ", "한화고딕 EL", 12, False, True),
    "summary2": ("• ", "한화고딕 EL", 12, False, True),
    "insight":  ("➔ ", "한화고딕 B",  12, True,  True),
}
DEFAULT_STYLE = ("", "한화고딕 EL", 12, False, True)


def parse_sections(text: str):
    matches = list(TAG_RE.finditer(text))
    return [
        (m.group(1).lower(), text[m.end():matches[i+1].start() if i+1 < len(matches) else len(text)].strip())
        for i, m in enumerate(matches)
        if text[m.end():matches[i+1].start() if i+1 < len(matches) else len(text)].strip()
    ]


def find_shape_by_index(prs: Presentation, shape_index: int, slide_index: int = 0):
    """특정 슬라이드의 특정 인덱스 shape 반환"""
    if slide_index >= len(prs.slides):
        return None, None
    
    slide = prs.slides[slide_index]
    shapes = list(slide.shapes)
    
    if shape_index >= len(shapes):
        return None, None
    
    return slide, shapes[shape_index]


def list_all_shapes(pptx_path: str):
    """디버깅용: 모든 슬라이드의 shape 정보 출력"""
    prs = Presentation(pptx_path)
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n=== 슬라이드 {slide_idx} ===")
        for i, shape in enumerate(slide.shapes):
            name = getattr(shape, "name", "N/A")
            has_tf = hasattr(shape, "has_text_frame") and shape.has_text_frame
            text_preview = ""
            if has_tf and shape.text_frame.text:
                text_preview = shape.text_frame.text[:30].replace('\n', ' ') + "..."
            print(f"  [{i}] {name} (text_frame: {has_tf}) {text_preview}")


def add_styled_run(paragraph, text, font_name, font_size, underline=False, color=None):
    r = paragraph.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.underline = underline
    if color:
        r.font.color.rgb = color


def set_number_and_date(prs: Presentation, number: str, date: str, 
                        shape_index: int = 4, slide_index: int = 0):
    """숫자와 날짜를 특정 TextBox에 입력"""
    _, shape = find_shape_by_index(prs, shape_index, slide_index)
    
    if not shape:
        raise ValueError(f'슬라이드 {slide_index}의 {shape_index}번째 shape을 찾지 못했습니다.')
    if not shape.has_text_frame:
        raise ValueError(f'{shape_index}번째 shape에 text_frame이 없습니다.')

    tf = shape.text_frame
    tf.clear()
    
    # 텍스트 구성 (예: "제25호 | 2025년 12월 26일")
    combined_text = f"제{number}호 | {date}"
    
    p = tf.paragraphs[0]
    add_styled_run(p, combined_text, "한화고딕 L", 11, color=RGBColor(0x6C, 0x6A, 0x67))


def set_textbox_from_summarizedtxt(prs: Presentation, text: str, 
                                    shape_index: int = 15, slide_index: int = 0):
    """특정 슬라이드의 특정 인덱스 shape에 텍스트 삽입"""
    _, shape = find_shape_by_index(prs, shape_index, slide_index)
    
    if not shape:
        raise ValueError(f'슬라이드 {slide_index}의 {shape_index}번째 shape을 찾지 못했습니다.')
    if not shape.has_text_frame:
        raise ValueError(f'{shape_index}번째 shape에 text_frame이 없습니다.')

    tf = shape.text_frame
    tf.clear()
    sections = parse_sections(text)

    if not sections:
        add_styled_run(tf.paragraphs[0], text.strip(), "한화고딕 EL", 12)
        return

    first_para_used = False
    for tag, content in sections:
        prefix, font_name, font_size, underline, split = TAG_STYLES.get(tag, DEFAULT_STYLE)
        lines = [ln.strip() for ln in content.splitlines() if ln.strip()] if split else [content.strip()]

        for line in filter(None, lines):
            p = tf.paragraphs[0] if not first_para_used and not tf.paragraphs[0].text else tf.add_paragraph()
            first_para_used = True
            add_styled_run(p, f"{prefix}{line}" if prefix else line, font_name, font_size, underline)

        if tag == "insight":
            add_styled_run(tf.add_paragraph(), " ", "한화고딕 EL", 9)


def create_report(pptx_in: str, pptx_out: str, number: str, date: str, 
                  text1: str, text2: str):
    """PPT 보고서 생성 (숫자/날짜 + 요약 텍스트 2개)"""
    prs = Presentation(pptx_in)
    
    # 1단계: 숫자와 날짜를 인덱스 4번 shape에 입력
    set_number_and_date(prs, number, date, shape_index=4, slide_index=0)
    
    # 2단계: 첫 번째 요약 텍스트를 인덱스 15번 shape에 입력
    set_textbox_from_summarizedtxt(prs, text1, shape_index=15, slide_index=0)
    
    # 3단계: 두 번째 요약 텍스트를 인덱스 16번 shape에 입력
    set_textbox_from_summarizedtxt(prs, text2, shape_index=16, slide_index=0)
    
    # 저장
    prs.save(pptx_out)


if __name__ == "__main__":
    
    # 사용자 입력 받기
    number = input("리포트 발행 호수를 입력하세요 (예: 25): ")
    date = input("리포트 발행 날짜를 입력하세요 (예: 2025년 12월 26일): ")
    
    # 보고서 생성
    create_report(
        pptx_in="AIWeeklyReport_format.pptx",
        pptx_out="output.pptx",
        number=number,
        date=date,
        text1=summarized_text,
        text2=summarized_text2
    )
    
    print("완료! output.pptx가 생성되었습니다.")